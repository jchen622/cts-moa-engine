#!/usr/bin/env python3
"""Build the explainer deck: how the MOA sourcing engine works.

Audience: the CTS AE small team and anyone the tool might be handed to. Every
number and example on these slides came from a real run, not an illustration.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

INK       = RGBColor(0x1C, 0x22, 0x2B)
INK_SOFT  = RGBColor(0x53, 0x5C, 0x69)
ACCENT    = RGBColor(0x1F, 0x4E, 0x79)
ACCENT_LT = RGBColor(0x9D, 0xB6, 0xCE)
GREEN     = RGBColor(0x1E, 0x6B, 0x52)
RED       = RGBColor(0x9B, 0x3A, 0x3A)
BOX_FILL  = RGBColor(0xF3, 0xF5, 0xF8)
BOX_LINE  = RGBColor(0xB8, 0xC2, 0xCE)
MUTED     = RGBColor(0x76, 0x80, 0x8C)
RULE      = RGBColor(0xD8, 0xDD, 0xE3)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG   = RGBColor(0x2B, 0x31, 0x3B)
CODE_FG   = RGBColor(0xE8, 0xEC, 0xF1)
TBL_HDR   = RGBColor(0xE8, 0xED, 0xF3)

FONT, MONO = "Arial", "Menlo"
SW, SH = Inches(13.333), Inches(7.5)
M = Inches(0.62)
CW = SW - 2 * M


def deck():
    p = Presentation()
    p.slide_width, p.slide_height = SW, SH
    return p


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def tb(s, x, y, w, h, wrap=True):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, size=12, bold=False, color=INK, after=4, before=0,
         italic=False, first=False, align=PP_ALIGN.LEFT, line=None, font=FONT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after, p.space_before = Pt(after), Pt(before)
    if line:
        p.line_spacing = line
    r = p.add_run()
    r.text = text
    r.font.name, r.font.size, r.font.bold, r.font.italic = font, Pt(size), bold, italic
    r.font.color.rgb = color
    return p


def rule(s, x, y, w, color=RULE, h=Emu(9525)):
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    ln.fill.solid(); ln.fill.fore_color.rgb = color
    ln.line.fill.background(); ln.shadow.inherit = False
    return ln


def header(s, kicker, title):
    t = tb(s, M, Inches(0.36), CW, Inches(0.26))
    para(t, kicker.upper(), size=10, bold=True, color=ACCENT, first=True, after=0)
    t2 = tb(s, M, Inches(0.63), CW, Inches(0.44))
    para(t2, title, size=24, bold=True, color=INK, first=True, after=0)
    rule(s, M, Inches(1.17), CW)


def footer(s, n, note=""):
    tf = tb(s, M, SH - Inches(0.42), CW - Inches(0.7), Inches(0.22))
    para(tf, note or "CTS MOA sourcing engine — how it works", size=8,
         color=MUTED, first=True, after=0)
    t2 = tb(s, SW - M - Inches(0.5), SH - Inches(0.42), Inches(0.5), Inches(0.22))
    para(t2, str(n), size=8, color=MUTED, first=True, after=0, align=PP_ALIGN.RIGHT)


def panel(s, x, y, w, h, fill=WHITE, line=RULE, rounded=False):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    b = s.shapes.add_shape(shape, x, y, w, h)
    if rounded:
        b.adjustments[0] = 0.05
    b.fill.solid(); b.fill.fore_color.rgb = fill
    b.line.color.rgb = line; b.line.width = Pt(0.75)
    b.shadow.inherit = False; b.text_frame.text = ""
    return b


def code(s, x, y, w, lines, size=10.5, pad=Inches(0.16)):
    h = pad * 2 + Inches(0.205) * len(lines)
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    b.fill.solid(); b.fill.fore_color.rgb = CODE_BG
    b.line.fill.background(); b.shadow.inherit = False; b.text_frame.text = ""
    tf = tb(s, x + pad, y + pad, w - pad * 2, h - pad * 2)
    for i, ln in enumerate(lines):
        para(tf, ln, size=size, color=CODE_FG, first=(i == 0), after=1,
             font=MONO, line=1.0)
    return h


def bullets(s, x, y, w, h, items, size=11.5, gap=6, marker="— "):
    tf = tb(s, x, y, w, h)
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            lead, rest = it
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(gap); p.line_spacing = 1.06
            r = p.add_run(); r.text = marker + lead
            r.font.name, r.font.size, r.font.bold = FONT, Pt(size), True
            r.font.color.rgb = INK
            r2 = p.add_run(); r2.text = rest
            r2.font.name, r2.font.size = FONT, Pt(size)
            r2.font.color.rgb = INK_SOFT
        else:
            para(tf, marker + it, size=size, color=INK_SOFT, first=(i == 0),
                 after=gap, line=1.06)
    return tf


def table(s, x, y, w, rows, col_w, size=9.5, row_h=Inches(0.27),
          hdr_h=Inches(0.29), hdr_color=ACCENT):
    nr, nc = len(rows), len(rows[0])
    t = s.shapes.add_table(nr, nc, x, y, w, hdr_h + row_h * (nr - 1)).table
    t.first_row = True; t.horz_banding = False
    for j, cwid in enumerate(col_w):
        t.columns[j].width = cwid
    t.rows[0].height = hdr_h
    for i in range(1, nr):
        t.rows[i].height = row_h
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = t.cell(i, j)
            c.margin_left = c.margin_right = Inches(0.06)
            c.margin_top = c.margin_bottom = Inches(0.02)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid(); c.fill.fore_color.rgb = TBL_HDR if i == 0 else WHITE
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.space_before = Pt(0); p.space_after = Pt(0)
            r = p.add_run(); r.text = str(val) if val else " "
            r.font.name = FONT; r.font.size = Pt(size)
            r.font.bold = (i == 0)
            r.font.color.rgb = hdr_color if i == 0 else INK
    return t


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def chip(s, x, y, w, h, text, fill, txt_color=WHITE, size=10):
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    b.adjustments[0] = 0.18
    b.fill.solid(); b.fill.fore_color.rgb = fill
    b.line.fill.background(); b.shadow.inherit = False
    tf = b.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.06)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(0)
    r = p.add_run(); r.text = text
    r.font.name, r.font.size, r.font.bold = FONT, Pt(size), True
    r.font.color.rgb = txt_color
    return b


# ---------------------------------------------------------------- slides
def s1(prs):
    s = blank(prs)
    rule(s, M, Inches(1.5), Inches(1.4), color=ACCENT, h=Inches(0.045))
    tf = tb(s, M, Inches(1.8), Inches(8.8), Inches(1.7))
    para(tf, "The MOA sourcing engine", size=38, bold=True, first=True, after=6)
    para(tf, "How new drug approvals become mini-review candidates — automatically",
         size=17, color=ACCENT, after=0)

    tf2 = tb(s, M, Inches(3.5), Inches(7.4), Inches(1.4))
    para(tf2, "A small program that watches every FDA approval, keeps only the "
              "genuinely novel agents, ranks them for the MOA mini-review series, and "
              "then names the clinical pharmacologist at the company who actually "
              "worked on each drug.",
         size=13, color=INK_SOFT, first=True, after=0, line=1.2)

    for i, (t, d) in enumerate([
            ("Runs when you run it", "no background surprises; start and stop at will"),
            ("Writes nothing by default", "every write needs an explicit --go"),
            ("Never sends email", "it drafts invitations; you send them")]):
        y = Inches(5.0) + Inches(0.62) * i
        chip(s, M, y, Inches(0.3), Inches(0.3), "✓", GREEN, size=12)
        t3 = tb(s, M + Inches(0.45), y + Inches(0.02), Inches(6.8), Inches(0.5))
        p = t3.paragraphs[0]; p.space_after = Pt(0)
        r = p.add_run(); r.text = t + " — "
        r.font.name, r.font.size, r.font.bold = FONT, Pt(12), True
        r.font.color.rgb = INK
        r2 = p.add_run(); r2.text = d
        r2.font.name, r2.font.size = FONT, Pt(12); r2.font.color.rgb = INK_SOFT

    p = panel(s, M + Inches(8.1), Inches(1.8), CW - Inches(8.1), Inches(4.9),
              fill=BOX_FILL, line=BOX_LINE, rounded=True)
    tf4 = tb(s, M + Inches(8.35), Inches(2.05), CW - Inches(8.6), Inches(4.4))
    para(tf4, "IN ONE RUN, TODAY", size=9.5, bold=True, color=ACCENT, first=True, after=10)
    for big, small in [("29,278", "FDA applications indexed"),
                       ("27", "novel agents in the current queue"),
                       ("24", "with a named clinical pharmacologist to write to"),
                       ("18/19", "published MOA papers the filter would have found")]:
        para(tf4, big, size=25, bold=True, color=ACCENT, after=0)
        para(tf4, small, size=10, color=INK_SOFT, after=13, line=1.05)
    footer(s, 1)
    notes(s, "Opening. Everything quoted in this deck comes from a real run on "
             "21 Aug 2026, not a mock-up. The three guarantees on the left are the "
             "design constraints the tool was built to honour.")


def s2(prs):
    s = blank(prs)
    header(s, "why", "The problem: the series has never held a steady rate")

    cx, cy, cwid = M, Inches(1.5), Inches(5.6)
    para(tb(s, cx, cy, cwid, Inches(0.3)), "MOA MINI-REVIEWS PUBLISHED PER YEAR",
         size=9.5, bold=True, color=ACCENT, first=True, after=0)
    data = [("2023", 1, ""), ("2024", 11, ""), ("2025", 2, ""), ("2026", 5, " (to Aug)")]
    top, bh, gap, lw = cy + Inches(0.42), Inches(0.36), Inches(0.54), Inches(0.6)
    maxw = cwid - lw - Inches(1.3)
    for i, (yr, n, sfx) in enumerate(data):
        y = top + gap * i
        para(tb(s, cx, y + Inches(0.05), lw, Inches(0.3)), yr, size=11.5, bold=True,
             first=True, after=0)
        bw = int(maxw * n / 11)
        b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx + lw, y, bw, bh)
        b.fill.solid()
        b.fill.fore_color.rgb = ACCENT if n == 11 else (RED if n <= 2 else ACCENT_LT)
        b.line.fill.background(); b.shadow.inherit = False; b.text_frame.text = ""
        para(tb(s, cx + lw + bw + Inches(0.1), y + Inches(0.06), Inches(1.5),
                Inches(0.3), wrap=False),
             f"{n}{sfx}", size=11, bold=True, first=True, after=0)

    para(tb(s, cx, Inches(4.05), cwid, Inches(0.9)),
         "A launch batch, a collapse, a partial recovery. The 2025 trough is what "
         "happens when every step of sourcing needs a person to remember to start it.",
         size=11.5, italic=True, color=INK_SOFT, first=True, after=0, line=1.15)

    rx = M + Inches(6.1)
    rw = CW - Inches(6.1)
    para(tb(s, rx, Inches(1.5), rw, Inches(0.3)), "THE PROCESS THIS REPLACES",
         size=9.5, bold=True, color=ACCENT, first=True, after=0)
    steps = [("Download the Purple Book by hand",
              "12 monthly files sat in Drive, unread"),
             ("Skim for anything new", "no record of what had been reviewed"),
             ("Look up who to contact",
              "a 44-company grid with one row filled in"),
             ("Send an invitation, sometimes", "no queue, no owner, no follow-up")]
    for i, (t, d) in enumerate(steps):
        y = Inches(1.86) + Inches(0.82) * i
        panel(s, rx, y, rw, Inches(0.68), fill=BOX_FILL, line=BOX_LINE)
        tf = tb(s, rx + Inches(0.16), y + Inches(0.1), rw - Inches(0.32), Inches(0.5))
        para(tf, f"{i+1}.  {t}", size=11, bold=True, first=True, after=2)
        para(tf, d, size=9.5, italic=True, color=MUTED, after=0)
        if i < 3:
            para(tb(s, rx + rw / 2, y + Inches(0.68), Inches(0.3), Inches(0.14)),
                 "▼", size=8, color=ACCENT_LT, first=True, after=0)

    para(tb(s, rx, Inches(5.35), rw, Inches(0.5)),
         "Every arrow needs a human to push it. Miss one month and the queue is empty.",
         size=10.5, italic=True, color=RED, first=True, after=0, line=1.15)
    footer(s, 2)
    notes(s, "Counts are from PubMed, filtered to the series title convention, and match "
             "the team deck. The right-hand column is not a caricature: the twelve "
             "purplebook-search downloads and the mostly-empty contact grid are both "
             "sitting in the ASCPT Drive folder today.")


def s3(prs):
    s = blank(prs)
    header(s, "how it works", "The pipeline, end to end")

    stages = [
        ("1. COLLECT", "Two FDA feeds", ["Drugs@FDA bulk files", "Purple Book (CBER)"], ACCENT),
        ("2. FILTER", "Is it a novel agent?", ["known-moiety index", "new-route rule"], RED),
        ("3. CLASSIFY", "What kind of drug?", ["modality from INN stem", "coverage-gap flag"], ACCENT_LT),
        ("4. ENRICH", "Who worked on it?", ["Phase 1 clin pharm papers", "author affiliations"], ACCENT),
        ("5. RANK", "How urgent?", ["0–100 score", "reasons recorded"], ACCENT_LT),
        ("6. DELIVER", "Where it lands", ["queue sheet", "outreach list + drafts"], GREEN),
    ]
    bw2 = Inches(1.94)
    gapx = Inches(0.12)
    for i, (num, title, items, col) in enumerate(stages):
        x = M + (bw2 + gapx) * i
        panel(s, x, Inches(1.5), bw2, Inches(2.05), fill=WHITE)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.5), bw2, Inches(0.055))
        bar.fill.solid(); bar.fill.fore_color.rgb = col
        bar.line.fill.background(); bar.shadow.inherit = False; bar.text_frame.text = ""
        tf = tb(s, x + Inches(0.12), Inches(1.68), bw2 - Inches(0.24), Inches(1.8))
        para(tf, num, size=9, bold=True, color=col, first=True, after=4)
        para(tf, title, size=11.5, bold=True, color=INK, after=7, line=1.05)
        for it in items:
            para(tf, "· " + it, size=9.5, color=INK_SOFT, after=3, line=1.05)
        if i < len(stages) - 1:
            para(tb(s, x + bw2 - Inches(0.02), Inches(2.34), Inches(0.32),
                    Inches(0.3), wrap=False), "▸", size=11, color=ACCENT_LT,
                 first=True, after=0, align=PP_ALIGN.CENTER)

    para(tb(s, M, Inches(3.78), CW, Inches(0.3)),
         "WHAT COMES OUT — three files in one folder on your computer",
         size=9.5, bold=True, color=ACCENT, first=True, after=0)
    outs = [("MOA candidate queue",
             "Append-only. The machine fills the facts; the AE owner and Status "
             "columns are yours and are never overwritten."),
            ("MOA author outreach list",
             "Ranked candidates, each with the clinical pharmacologists to contact, "
             "why they were picked and the PubMed IDs behind it."),
            ("Invitation drafts",
             "One editable draft per candidate, opened in a browser or Word. Nothing is "
             "sent; you review, edit and send each one yourself.")]
    ow = Inches(3.95)
    for i, (t, d) in enumerate(outs):
        x = M + (ow + Inches(0.13)) * i
        panel(s, x, Inches(4.1), ow, Inches(1.35), fill=BOX_FILL, line=BOX_LINE,
              rounded=True)
        tf = tb(s, x + Inches(0.16), Inches(4.26), ow - Inches(0.32), Inches(1.1))
        para(tf, t, size=11.5, bold=True, color=ACCENT, first=True, after=5)
        para(tf, d, size=9.5, color=INK_SOFT, after=0, line=1.1)

    para(tb(s, M, Inches(5.72), CW, Inches(0.3)), "TWO CADENCES",
         size=9.5, bold=True, color=ACCENT, first=True, after=0)
    tf = tb(s, M, Inches(6.02), CW, Inches(0.8))
    para(tf, "Monthly — a light run tops up the queue so nothing is lost.  ·  "
             "Whenever you are ready to write — build the outreach list and the "
             "letters. Before ASCPT, drop in the programme or attendee list and the "
             "same people gain a poster number and a time.",
         size=11.5, color=INK_SOFT, first=True, after=0, line=1.15)
    footer(s, 3)
    notes(s, "Walk left to right. The key point for a sceptical listener is stage 2: "
             "collecting approvals is easy, deciding which are genuinely novel is the "
             "part that needed real work.")


def s4(prs):
    s = blank(prs)
    header(s, "step 1 — collect", "Where the data comes from, and why it takes two sources")

    lw = Inches(6.15)
    panel(s, M, Inches(1.48), lw, Inches(2.5), fill=WHITE)
    tf = tb(s, M + Inches(0.2), Inches(1.64), lw - Inches(0.4), Inches(2.2))
    para(tf, "SOURCE A — Drugs@FDA bulk files", size=12, bold=True, color=ACCENT,
         first=True, after=6)
    para(tf, "A 6 MB zip of 12 relational tables, rebuilt daily. Joining "
             "Applications × Submissions × Products gives every original approval and "
             "its FDA submission class.", size=10.5, color=INK_SOFT, after=8, line=1.12)
    para(tf, "Covers:  all NDAs, plus CDER-regulated biologics (BLA 761xxx)",
         size=10.5, after=4)
    para(tf, "Misses:  CBER products — cell and gene therapy, vaccines, plasma",
         size=10.5, color=RED, after=0)

    panel(s, M, Inches(4.12), lw, Inches(2.5), fill=WHITE)
    tf = tb(s, M + Inches(0.2), Inches(4.28), lw - Inches(0.4), Inches(2.2))
    para(tf, "SOURCE B — Purple Book monthly report", size=12, bold=True, color=ACCENT,
         first=True, after=6)
    para(tf, "Already a change report, not a snapshot: an N/R/U column marks each row "
             "New, added in Release, or Updated. No month-to-month diffing needed — "
             "which the manual process never exploited.",
         size=10.5, color=INK_SOFT, after=8, line=1.12)
    para(tf, "Covers:  CBER biologics — exactly the modality gaps the series has",
         size=10.5, after=4)
    para(tf, "Excludes:  351(k) biosimilars and interchangeables, by rule",
         size=10.5, color=INK_SOFT, after=0)

    rx = M + Inches(6.45)
    rw = CW - Inches(6.45)
    panel(s, rx, Inches(1.48), rw, Inches(1.75), fill=BOX_FILL, line=BOX_LINE,
          rounded=True)
    tf = tb(s, rx + Inches(0.2), Inches(1.64), rw - Inches(0.4), Inches(1.45))
    para(tf, "WHY BOTH — the proof", size=9.5, bold=True, color=ACCENT, first=True, after=6)
    para(tf, "BLA 125730 (StrataGraft, a cell therapy) is absent from Drugs@FDA "
             "entirely. Only 73 of the classic 125xxx CBER applications appear there "
             "at all. Without the Purple Book, the series would never see a cell or "
             "gene therapy approval.", size=10.5, color=INK, after=0, line=1.12)

    panel(s, rx, Inches(3.38), rw, Inches(3.24), fill=WHITE, line=RED)
    tf = tb(s, rx + Inches(0.2), Inches(3.54), rw - Inches(0.4), Inches(2.9))
    para(tf, "⚠  A TRAP WORTH KNOWING", size=9.5, bold=True, color=RED, first=True, after=6)
    para(tf, "The obvious approach — the openFDA web API — silently returns wrong "
             "answers here. It matches conditions across the whole application "
             "document rather than within one submission, so:",
         size=10.5, color=INK_SOFT, after=7, line=1.12)
    code(s, rx + Inches(0.2), Inches(4.68), rw - Inches(0.4),
         ['search=date:[2026] AND class:"TYPE 1"'], size=9)
    tf2 = tb(s, rx + Inches(0.2), Inches(5.22), rw - Inches(0.4), Inches(1.2))
    para(tf2, "claims 279 results and returns approvals from 2008, 2014, 2016, 2019 "
              "and 2020. The bulk files are used instead precisely because they allow "
              "a correct per-submission join.",
         size=10.5, color=INK, first=True, after=0, line=1.12)
    footer(s, 4)
    notes(s, "If someone asks 'why not just use the FDA API?', this is the slide. The "
             "openFDA behaviour is not a bug in our code — it is how their nested "
             "search works — but it would have produced a queue full of decades-old "
             "drugs with nobody noticing.")


def s5(prs):
    s = blank(prs)
    header(s, "step 2 — filter", "What counts as a “novel agent” — the heart of the tool")

    para(tb(s, M, Inches(1.42), CW, Inches(0.4)),
         "Novelty is decided against a known-moiety index built from all 29,278 "
         "applications in the FDA record — not a keyword list. Every verdict carries a "
         "plain-text reason you can argue with.",
         size=11.5, color=INK_SOFT, first=True, after=0, line=1.12)

    lw = Inches(6.15)
    panel(s, M, Inches(1.95), lw, Inches(2.75), fill=WHITE, line=GREEN)
    chip(s, M + Inches(0.2), Inches(2.1), Inches(1.15), Inches(0.28), "INCLUDED", GREEN, size=9)
    bullets(s, M + Inches(0.2), Inches(2.5), lw - Inches(0.4), Inches(2.1), [
        ("First approval of a moiety ", "anywhere in the FDA record"),
        ("A combination with one new component ", "— cefepime + zidebactam"),
        ("New active ingredient ", "— an enantiomer or salt (esketamine)"),
        ("Known molecule by a new route ", "— intrathecal Zolgensma (ITVISMA), "
         "ophthalmic bevacizumab (LYTENAVA)"),
    ], size=10.5, gap=7)

    panel(s, M, Inches(4.85), lw, Inches(2.0), fill=WHITE, line=RED)
    chip(s, M + Inches(0.2), Inches(5.0), Inches(1.15), Inches(0.28), "EXCLUDED", RED, size=9)
    bullets(s, M + Inches(0.2), Inches(5.4), lw - Inches(0.4), Inches(1.4), [
        "Generics (ANDAs), biosimilars and interchangeables",
        "Reformulations and new combinations of old drugs",
        "Repeat applications by the same route — a reformulated COVID-19 vaccine",
        "Plasma-derived products — immune globulin, prothrombin complex",
    ], size=10.5, gap=5)

    rx = M + Inches(6.45)
    rw = CW - Inches(6.45)
    para(tb(s, rx, Inches(1.95), rw, Inches(0.3)),
         "REAL VERDICTS FROM THE LAST RUN", size=9.5, bold=True, color=ACCENT,
         first=True, after=0)
    rows = [["Agent", "Verdict"],
            ["orforglipron", "KEEP — first approval of this moiety"],
            ["zidebactam\n(with cefepime)", "KEEP — novel component in a combination"],
            ["onasemnogene\nabeparvovec", "KEEP — new route: intrathecal; new brand ITVISMA"],
            ["bevacizumab", "KEEP — new brand LYTENAVA on an NME-classed original"],
            ["COVID-19 vaccine\nmRNA", "DROP — previously approved moiety and route"],
            ["ranibizumab", "DROP — previously approved moiety and route"],
            ["immune globulin", "DROP — plasma-derived, non-novel class"]]
    table(s, rx, Inches(2.28), rw, rows, [Inches(1.65), rw - Inches(1.65)],
          size=9.5, row_h=Inches(0.52))
    para(tb(s, rx, Inches(6.28), rw, Inches(0.5)),
         "The new-route rule is why the two “known molecule” rows are keeps, not drops.",
         size=9.5, italic=True, color=MUTED, first=True, after=0, line=1.1)
    footer(s, 5)
    notes(s, "This is the slide to linger on. The distinction that matters: FDA's NME "
             "class alone is not enough — it would have let through a new BLA for "
             "bevacizumab as if it were a new molecule. The moiety index catches that; "
             "the new-route rule then decides it is interesting anyway, for a different "
             "and explicit reason.")


def s6(prs):
    s = blank(prs)
    header(s, "steps 3–5", "Classify, enrich, and rank")

    lw = Inches(6.15)
    panel(s, M, Inches(1.45), lw, Inches(1.62), fill=WHITE)
    tf = tb(s, M + Inches(0.2), Inches(1.6), lw - Inches(0.4), Inches(1.35))
    para(tf, "CLASSIFY — modality from the drug's name", size=11.5, bold=True,
         color=ACCENT, first=True, after=6)
    para(tf, "INN stems are a naming convention, so the name itself encodes the "
             "modality. The list follows the 2021 WHO revision that replaced “-mab” "
             "with -tug, -bart, -mig and -ment — without it, current antibodies such "
             "as veligrotug get filed as small molecules.",
         size=10.5, color=INK_SOFT, after=0, line=1.12)

    panel(s, M, Inches(3.2), lw, Inches(1.62), fill=WHITE)
    tf = tb(s, M + Inches(0.2), Inches(3.35), lw - Inches(0.4), Inches(1.35))
    para(tf, "ENRICH — ask PubMed two questions", size=11.5, bold=True, color=ACCENT,
         first=True, after=6)
    para(tf, "1.  Has a mechanism-of-action review already been written? If so, don't "
             "invite a duplicate.\n"
             "2.  Who publishes clinical pharmacology on this drug? Those names become "
             "the suggested authors.",
         size=10.5, color=INK_SOFT, after=0, line=1.15)

    panel(s, M, Inches(4.95), lw, Inches(1.75), fill=BOX_FILL, line=BOX_LINE, rounded=True)
    tf = tb(s, M + Inches(0.2), Inches(5.1), lw - Inches(0.4), Inches(1.5))
    para(tf, "GAP FLAGS", size=9.5, bold=True, color=ACCENT, first=True, after=6)
    para(tf, "A candidate that fills a hole in the series scores higher. The gap list "
             "comes straight from the team deck: cell & gene therapy, CAR-T, "
             "radioligand therapy, siRNA/ASO, incretins & cardiometabolic, vaccines, "
             "AI-derived assets.", size=10.5, color=INK, after=0, line=1.12)

    rx = M + Inches(6.45)
    rw = CW - Inches(6.45)
    para(tb(s, rx, Inches(1.45), rw, Inches(0.3)),
         "RANK — a legible 0–100 score", size=11.5, bold=True, color=ACCENT,
         first=True, after=0)
    para(tb(s, rx, Inches(1.78), rw, Inches(0.4)),
         "Not a black box. Every term is recorded on the row, so the group can "
         "disagree with the weights and change them in one place.",
         size=10.5, color=INK_SOFT, first=True, after=0, line=1.12)

    panel(s, rx, Inches(2.35), rw, Inches(2.45), fill=WHITE, line=ACCENT)
    tf = tb(s, rx + Inches(0.2), Inches(2.5), rw - Inches(0.4), Inches(2.2))
    para(tf, "WORKED EXAMPLE — insulin icodec", size=9.5, bold=True, color=ACCENT,
         first=True, after=8)
    for term, pts, col in [("first-in-human moiety", "+40", INK),
                           ("fills coverage gap: incretins & cardiometabolic", "+22", INK),
                           ("no existing MOA review found", "+8", INK),
                           ("clin pharm literature exists, 13 papers", "+6", INK),
                           ("approved this year", "+8", INK)]:
        p = tf.add_paragraph(); p.space_after = Pt(5); p.line_spacing = 1.05
        r = p.add_run(); r.text = f"{pts:>5}   "
        r.font.name, r.font.size, r.font.bold = MONO, Pt(10), True
        r.font.color.rgb = ACCENT
        r2 = p.add_run(); r2.text = term
        r2.font.name, r2.font.size = FONT, Pt(10); r2.font.color.rgb = col
    p = tf.add_paragraph(); p.space_before = Pt(6); p.space_after = Pt(0)
    r = p.add_run(); r.text = "   84   "
    r.font.name, r.font.size, r.font.bold = MONO, Pt(13), True
    r.font.color.rgb = ACCENT
    r2 = p.add_run(); r2.text = "total — top of the current queue"
    r2.font.name, r2.font.size, r2.font.bold = FONT, Pt(11), True
    r2.font.color.rgb = INK

    panel(s, rx, Inches(4.95), rw, Inches(1.75), fill=WHITE, line=RULE)
    tf = tb(s, rx + Inches(0.2), Inches(5.1), rw - Inches(0.4), Inches(1.5))
    para(tf, "SCORES GO DOWN TOO", size=9.5, bold=True, color=ACCENT, first=True, after=6)
    for t in ["−30  a MOA review already exists",
              "−15  diagnostic rather than therapeutic",
              "−8   no clin pharm literature, authors unclear"]:
        para(tf, t, size=10, color=INK_SOFT, after=4, font=MONO, line=1.05)
    footer(s, 6)
    notes(s, "The scoring is deliberately arithmetic rather than clever. If the group "
             "thinks a coverage gap should outweigh first-in-class, that is a one-line "
             "change in classify.py and the deck's numbers move with it.")


def s7(prs):
    s = blank(prs)
    header(s, "the payoff", "Who to write to — by name, with the evidence")

    para(tb(s, M, Inches(1.40), CW, Inches(0.40)),
         "The objective is reaching the clinical pharmacologist who actually worked on "
         "the drug. The tool finds them in PubMed and says why it picked them.",
         size=11.5, color=INK_SOFT, first=True, after=0, line=1.12)

    para(tb(s, M, Inches(1.90), CW, Inches(0.26)),
         "THREE PASSES — because author position means different things in different papers",
         size=9.5, bold=True, color=ACCENT, first=True, after=0)
    rows = [["", "What it searches", "Where the clinical pharmacologist sits"],
            ["1", "Food effect · DDI · relative bioavailability · organ impairment · "
                  "mass balance / ADME · thorough QT · healthy volunteers",
             "FIRST author — they designed and ran it"],
            ["2", "Population PK · exposure–response · PBPK",
             "First or senior author"],
            ["3", "Dose escalation / first-in-human  (fallback only)",
             "ANYWHERE — first and last are the treating clinicians"]]
    table(s, M, Inches(2.20), CW, rows,
          [Inches(0.34), Inches(6.5), CW - Inches(6.84)], size=9.5, row_h=Inches(0.46))

    lw = Inches(6.15)
    panel(s, M, Inches(4.00), lw, Inches(2.10), fill=WHITE, line=RULE)
    tf = tb(s, M + Inches(0.18), Inches(4.13), lw - Inches(0.36), Inches(1.85))
    para(tf, "THE CASE THAT PROVES IT ISN'T JUST POSITION", size=9, bold=True,
         color=ACCENT, first=True, after=6)
    para(tf, "Pivekimab sunirine, an AbbVie oncology ADC. The Phase 1 paper has "
             "24 authors.", size=10, color=INK_SOFT, after=6, line=1.1)
    for t, c in [("first    Naveen Pemmaraju — MD Anderson", RED),
                 ("20–22    Yining Du · Sribalaji Lakshmikanthan · "
                  "Jalaja Potluri — AbbVie", GREEN),
                 ("last     Naval G Daver — MD Anderson", RED)]:
        para(tf, t, size=9.5, color=c, after=3, line=1.05)
    para(tf, "The tool returns the three AbbVie authors and neither clinician. Pinned "
             "as a test so it cannot regress.",
         size=9.5, italic=True, color=MUTED, after=0, line=1.1)

    rx = M + Inches(6.45)
    rw = CW - Inches(6.45)
    panel(s, rx, Inches(4.00), rw, Inches(2.10), fill=BOX_FILL, line=BOX_LINE, rounded=True)
    tf = tb(s, rx + Inches(0.18), Inches(4.13), rw - Inches(0.36), Inches(1.85))
    para(tf, "TWO THINGS IT HAS TO GET RIGHT", size=9, bold=True, color=ACCENT,
         first=True, after=6)
    para(tf, "Acquisitions. Tebipenem is approved to GSK, but every paper comes out of "
             "Spero, whom GSK bought. Baxdrostat is AstraZeneca; the papers are CinCor. "
             "The tool infers the company that actually ran the programme.",
         size=9.5, color=INK, after=6, line=1.1)
    para(tf, "People move. PubMed records the affiliation at publication, so it reports "
             "where the work was done and adds “now at …”.",
         size=9.5, color=INK, after=0, line=1.1)

    para(tb(s, M, Inches(6.28), CW, Inches(0.36)),
         "Result on the live queue: 24 of 27 candidates get a named clinical "
         "pharmacologist, 20 of them at the sponsor or the company that ran the "
         "programme — against a baseline where most read NEEDS LOOKUP.",
         size=11, bold=True, color=ACCENT, first=True, after=0, line=1.12)
    footer(s, 7)
    notes(s, "This replaces the old version of this slide, which framed the payoff as "
             "the ASCPT meeting. The meeting is one channel for reaching the person; "
             "finding the person is the job.\n\n"
             "Every name here is real output. The pivekimab example is the important "
             "one: a naive 'take the first author' rule returns an MD Anderson "
             "oncologist, which is exactly the wrong person to invite to write a "
             "clinical pharmacology mini-review.\n\n"
             "The tool also harvests corresponding-author emails from the affiliation "
             "string when PubMed carries them.")


def s8(prs):
    s = blank(prs)
    header(s, "does it work?", "Backtest against the papers the series actually published")

    para(tb(s, M, Inches(1.42), CW, Inches(0.42)),
         "The 19 published MOA mini-reviews are the only labelled data that exists. The "
         "test: re-scan every FDA approval since 2015 and ask whether the filter would "
         "have surfaced the drugs that became papers.",
         size=11.5, color=INK_SOFT, first=True, after=0, line=1.12)

    panel(s, M, Inches(2.0), Inches(3.5), Inches(1.9), fill=WHITE, line=GREEN)
    tf = tb(s, M + Inches(0.25), Inches(2.2), Inches(3.0), Inches(1.6))
    para(tf, "RECALL", size=9.5, bold=True, color=ACCENT, first=True, after=4)
    para(tf, "18 / 19", size=34, bold=True, color=GREEN, after=4)
    para(tf, "published MOA drugs are surfaced by the filter",
         size=10.5, color=INK_SOFT, after=0, line=1.1)

    panel(s, M + Inches(3.75), Inches(2.0), CW - Inches(3.75), Inches(1.9),
          fill=WHITE, line=RULE)
    tf = tb(s, M + Inches(3.95), Inches(2.2), CW - Inches(4.15), Inches(1.6))
    para(tf, "THE ONE MISS — molnupiravir", size=11.5, bold=True, color=ACCENT,
         first=True, after=6)
    para(tf, "It only ever held an Emergency Use Authorization and never received full "
             "FDA approval, so it appears nowhere in Drugs@FDA. This is a limit of the "
             "source, not a defect in the filter — EUA products are invisible to the "
             "tool and have to be spotted by a person.",
         size=10.5, color=INK_SOFT, after=0, line=1.12)

    panel(s, M, Inches(4.08), CW, Inches(1.55), fill=BOX_FILL, line=BOX_LINE, rounded=True)
    tf = tb(s, M + Inches(0.25), Inches(4.26), CW - Inches(0.5), Inches(1.25))
    para(tf, "WHAT THE BACKTEST CHANGED", size=9.5, bold=True, color=ACCENT,
         first=True, after=7)
    para(tf, "The first version scored 17/19 and excluded esketamine. Investigating "
             "showed why: esketamine is FDA class Type 2 — a new active ingredient, the "
             "enantiomer of an approved drug — not a Type 1 new molecular entity. Since "
             "the series had published it, Type 2 was added as a lower-scoring tier. "
             "The test found a real gap in the rules, which is the point of having it.",
         size=11, color=INK, after=0, line=1.15)

    para(tb(s, M, Inches(5.85), CW, Inches(0.3)), "RE-RUN IT AFTER ANY CHANGE",
         size=9.5, bold=True, color=ACCENT, first=True, after=0)
    code(s, M, Inches(6.15), Inches(5.2), ["$ python3 backtest.py"], size=11)
    para(tb(s, M + Inches(5.5), Inches(6.28), CW - Inches(5.5), Inches(0.5)),
         "A filter that rejects drugs the series actually published is not ready to trust.",
         size=10.5, italic=True, color=INK_SOFT, first=True, after=0, line=1.1)
    footer(s, 8)
    notes(s, "Worth saying out loud: the backtest is the reason to believe any of this. "
             "It also caught a real bug — the Type 2 gap — which is a good answer to "
             "'how do you know it's right?'")


def s9(prs):
    s = blank(prs)
    header(s, "running it", "Commands, control, and the safety rails")

    para(tb(s, M, Inches(1.42), CW, Inches(0.3)),
         "EVERYDAY USE — most people just press the buttons; this is the same thing",
         size=9.5, bold=True, color=ACCENT, first=True, after=0)
    code(s, M, Inches(1.72), Inches(7.4), [
        "$ python3 moa_engine.py check           # verify everything; writes nothing",
        "$ python3 moa_engine.py scan --days 120 # preview candidates; writes nothing",
        "$ python3 moa_engine.py update --go     # add new candidates to the queue",
        "$ python3 moa_engine.py dossier --go    # build the outreach list",
        "$ python3 moa_engine.py invites --go    # draft the letters, once read",
    ], size=10.5)

    para(tb(s, M, Inches(3.28), CW, Inches(0.3)), "TURNING THE SCHEDULE ON AND OFF",
         size=9.5, bold=True, color=ACCENT, first=True, after=0)
    code(s, M, Inches(3.58), Inches(7.4), [
        "$ python3 moa_engine.py start           # monthly, 7th at 09:17",
        "$ python3 moa_engine.py status          # on or off?",
        "$ python3 moa_engine.py stop            # removes it completely",
    ], size=10.5)

    para(tb(s, M, Inches(4.85), Inches(7.4), Inches(0.9)),
         "It is off by default. If the Mac is asleep the job runs at next wake, and a "
         "missed month costs nothing: each run re-scans a trailing window and "
         "de-duplicates by key, so it self-heals. Running twice adds nothing the "
         "second time.",
         size=10.5, color=INK_SOFT, first=True, after=0, line=1.15)

    rx = M + Inches(7.7)
    rw = CW - Inches(7.7)
    panel(s, rx, Inches(1.42), rw, Inches(2.90), fill=WHITE, line=GREEN)
    tf = tb(s, rx + Inches(0.2), Inches(1.58), rw - Inches(0.4), Inches(2.25))
    para(tf, "THE SAFETY RAILS", size=9.5, bold=True, color=GREEN, first=True, after=8)
    for t, d in [("Runs on its own",
                  "no AI, no Claude, no account, no API key, no server — "
                  "Python 3 and a connection to the FDA and PubMed, nothing else"),
                 ("Dry run by default",
                  "nothing is written without --go"),
                 ("No email, ever",
                  "invitations are drafted for you to send"),
                 ("Your columns are yours",
                  "AE owner and Status are never overwritten"),
                 ("No guessed contacts",
                  "a weak name match reports NEEDS LOOKUP")]:
        p = tf.add_paragraph(); p.space_after = Pt(7); p.line_spacing = 1.08
        r = p.add_run(); r.text = "✓  " + t + " — "
        r.font.name, r.font.size, r.font.bold = FONT, Pt(9.5), True
        r.font.color.rgb = INK
        r2 = p.add_run(); r2.text = d
        r2.font.name, r2.font.size = FONT, Pt(9.5); r2.font.color.rgb = INK_SOFT

    panel(s, rx, Inches(4.46), rw, Inches(2.16), fill=BOX_FILL, line=BOX_LINE, rounded=True)
    tf = tb(s, rx + Inches(0.2), Inches(4.60), rw - Inches(0.4), Inches(1.90))
    para(tf, "HANDING IT TO SOMEONE ELSE — ONE FILE", size=9.5, bold=True,
         color=ACCENT, first=True, after=8)
    para(tf, "Send one file, ~93 KB. Both the app and the command line.",
         size=10.5, color=INK, after=6, line=1.15)
    para(tf, "SEND THIS/CTS MOA Engine.command     Mac\n"
             "SEND THIS/CTS MOA Engine.bat             Windows",
         size=10, bold=True, color=ACCENT, after=6, line=1.25)
    para(tf, "Loaded from memory; first run creates a working folder beside it.",
         size=10, color=INK_SOFT, after=6, line=1.12)
    para(tf, "No account, no credential, no API key, no server. The only requirement "
             "is Python 3.",
         size=10, color=INK_SOFT, after=0, line=1.12)
    footer(s, 9)
    notes(s, "The portability slide matters for the durability question: the tool is "
             "not tied to one laptop or one person, though someone still has to choose "
             "to run it. The only requirement is Python 3 — there is nothing to sign "
             "in to, and a browser window is the whole interface.")


def s10(prs):
    s = blank(prs)
    header(s, "honesty", "What it does not do — and what the group should decide")

    lw = Inches(6.15)
    para(tb(s, M, Inches(1.45), lw, Inches(0.3)), "KNOWN LIMITATIONS",
         size=9.5, bold=True, color=ACCENT, first=True, after=0)
    bullets(s, M, Inches(1.72), lw, Inches(3.5), [
        ("EUA products are invisible ",
         "— molnupiravir never entered Drugs@FDA"),
        ("Route data is imperfect ",
         "— Drugs@FDA has no route column; it is parsed out of the dosage form, which "
         "sometimes holds a presentation instead"),
        ("Names come from the literature, not a staff list ",
         "— someone who never published on the drug is invisible, and 3 of 27 "
         "candidates are too new to have any papers at all"),
        ("Affiliations are as at publication ",
         "— people move. The tool flags a likely move but cannot confirm a current "
         "employer. LinkedIn is not an option: its terms prohibit automated access "
         "and it refuses non-browser requests outright, so this stays a manual check"),
        ("The ASCPT member directory is behind a login ",
         "— it is never scraped. The tool writes a short check list of the specific "
         "names it wants an answer for, which you or ASCPT fill in"),
        ("Programme matches are leads, not facts ",
         "— “sponsor presenting” means someone from that company has a poster"),
        ("CBER depends on Purple Book lag ",
         "— only January was published for 2026 at the time of writing"),
    ], size=10, gap=6)

    panel(s, M, Inches(5.5), lw, Inches(1.20), fill=WHITE, line=RED)
    tf = tb(s, M + Inches(0.2), Inches(5.63), lw - Inches(0.4), Inches(1.0))
    para(tf, "THE BIGGEST ONE", size=9.5, bold=True, color=RED, first=True, after=6)
    para(tf, "It runs on one person's laptop. The tool is portable, but somebody still "
             "has to own it — otherwise this rebuilds the exact dependency that "
             "produced the 2025 trough.",
         size=11, color=INK, after=0, line=1.15)

    rx = M + Inches(6.45)
    rw = CW - Inches(6.45)
    para(tb(s, rx, Inches(1.45), rw, Inches(0.3)), "FOR THE GROUP TO DECIDE",
         size=9.5, bold=True, color=ACCENT, first=True, after=0)
    qs = ["Who owns the queue — a named AE, a rotation, or the A-team?",
          "What is the score threshold for actually inviting someone?",
          "Do we chase the contact grid, or route invitations through ASCPT?",
          "Should the outreach list be shared with the whole AE group?",
          "Does this stay on a laptop, or move somewhere the journal controls?"]
    for i, q in enumerate(qs):
        y = Inches(1.8) + Inches(0.92) * i
        panel(s, rx, y, rw, Inches(0.78), fill=BOX_FILL, line=BOX_LINE, rounded=True)
        tf = tb(s, rx + Inches(0.18), y + Inches(0.13), rw - Inches(0.36), Inches(0.6))
        para(tf, q, size=10.5, color=INK, first=True, after=0, line=1.12)
    footer(s, 10)
    notes(s, "End on the ownership question rather than on the technology. The tool "
             "removes the effort but not the accountability, and the group has to place "
             "that somewhere.")


def main():
    prs = deck()
    for fn in (s1, s2, s3, s4, s5, s6, s7, s8, s9, s10):
        fn(prs)
    out = ("/Users/chenj345/Desktop/Desktop/05 Code and Tools/cts-moa-engine/"
           "How the MOA sourcing engine works.pptx")
    prs.save(out)
    print("saved:", out)


if __name__ == "__main__":
    main()
